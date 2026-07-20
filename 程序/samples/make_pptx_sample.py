"""生成测试用 PowerPoint（标题页 + 要点页 + 表格 + 演讲者备注）。

用法（在 程序/ 目录下）：python samples/make_pptx_sample.py
输出：samples/sample_slides.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "sample_slides.pptx"


def build():
    prs = Presentation()

    # 幻灯片 1：标题
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Observing Robot Failures in Classrooms"
    s1.placeholders[1].text = "A study of productive failure with social robots"
    s1.notes_slide.notes_text_frame.text = (
        "Speaker note: emphasise the 135 students across six classes.")

    # 幻灯片 2：要点
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Key Findings"
    body = s2.placeholders[1].text_frame
    body.text = "Robot failure produced the largest conceptual gain."
    for line in ("Lower reported social pressure",
                 "Effect held after novelty subsided",
                 "Consistent across 3 instructional methods"):
        p = body.add_paragraph()
        p.text = line

    # 幻灯片 3：表格
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Results by Condition"
    rows, cols = 3, 3
    tbl = s3.shapes.add_table(rows, cols, Inches(1), Inches(1.8),
                              Inches(8), Inches(2)).table
    data = [["Condition", "Participants", "Outcome"],
            ["Robot failure", "46 students", "High improvement"],
            ["Direct instruction", "44 students", "Small improvement"]]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val
            tbl.cell(r, c).text_frame.paragraphs[0].runs[0].font.size = Pt(14)

    prs.save(str(OUT))
    print(f"已生成：{OUT}（{len(prs.slides._sldIdLst)} 张幻灯片，含表格与备注）")


if __name__ == "__main__":
    build()
