"""PowerPoint (.pptx) 翻译（方向一 A3）。

与 docx_translator 同一思路：按 run 就地替换文字，字体/颜色/位置/动画天然
保留。覆盖：幻灯片文本框、表格（含分组形状内的）、演讲者备注。

依赖 python-pptx（MIT 许可，可自由商用）。
"""
from __future__ import annotations

from typing import Callable, List, Optional


class PptxUnsupported(Exception):
    """缺少 python-pptx 组件。"""


def available() -> bool:
    try:
        import pptx  # noqa: F401
        return True
    except Exception:  # noqa: BLE001
        return False


class _Unit:
    """一个待译段落：文字取自 runs，回填也只动这些 runs。"""

    __slots__ = ("runs", "text")

    def __init__(self, runs):
        self.runs = runs
        self.text = "".join(r.text for r in runs)


def _worth_translating(text: str, target_code: str = "zh") -> bool:
    from .languages import looks_like
    t = (text or "").strip()
    if len(t) < 2:
        return False
    if looks_like(t, target_code):
        return False                       # 已是目标语（B5）
    return sum(c.isalpha() for c in t) >= 2


def _carrier(runs):
    """选字数最多的 run 承载译文（保住主导格式，与 docx 一致）。"""
    return max(runs, key=lambda r: len(r.text))


def _iter_text_frames(shapes):
    """递归遍历形状，产出所有 text_frame（含分组形状与表格单元格）。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_text_frames(sh.shapes)
            continue
        if sh.has_text_frame:
            yield sh.text_frame
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def _collect(prs, target_code: str) -> List[_Unit]:
    units: List[_Unit] = []
    frames = []
    for slide in prs.slides:
        frames.extend(_iter_text_frames(slide.shapes))
        if slide.has_notes_slide:                       # 演讲者备注
            frames.append(slide.notes_slide.notes_text_frame)
    for tf in frames:
        for para in tf.paragraphs:
            runs = [r for r in para.runs if r.text]
            if not runs:
                continue
            u = _Unit(runs)
            if _worth_translating(u.text, target_code):
                units.append(u)
    return units


def _set_text(unit: _Unit, text: str) -> None:
    keep = _carrier(unit.runs)
    keep.text = text
    for r in unit.runs:
        if r is not keep:
            r.text = ""


def translate_pptx(
    input_path: str,
    output_path: str,
    cfg,
    translator=None,
    glossary=None,
    mock: bool = False,
    progress: Optional[Callable[[str, float], None]] = None,
    should_cancel: Optional[Callable[[], bool]] = None,
) -> dict:
    if not available():
        raise PptxUnsupported(
            "缺少 PowerPoint 支持组件。请运行：pip install python-pptx")

    from pptx import Presentation

    from .glossary import Glossary
    from .pipeline import (CancelledError, _estimate_line, _maybe_pangu,
                           _translated, make_translator)

    def report(msg: str, frac: float):
        if progress:
            progress(msg, max(0.0, min(1.0, frac)))

    def check_cancel():
        if should_cancel and should_cancel():
            raise CancelledError("已取消。")

    check_cancel()
    report("正在解析 PowerPoint…", 0.03)
    prs = Presentation(input_path)
    target = getattr(cfg, "target_lang", "zh") or "zh"
    units = _collect(prs, target)
    texts = [u.text for u in units]
    n_slides = len(prs.slides._sldIdLst)
    report(f"共 {n_slides} 张幻灯片、{len(texts)} 个待译段落", 0.08)

    if glossary is None:
        glossary = Glossary.load(cfg.resolved_glossary_path())
    if translator is None:
        head = texts[0][:150] if texts else ""
        body = next((t for t in texts if len(t) > 200), "")
        ctx = head + (("\n摘要节选：" + body[:250]) if body else "")
        translator = make_translator(cfg, mock=mock, doc_context=ctx)
    if texts and not mock:
        try:
            report(_estimate_line(translator, texts, cfg), 0.09)
        except Exception:  # noqa: BLE001
            pass

    def tcb(done: int, total: int):
        check_cancel()
        report(f"正在翻译… 第 {done}/{total} 批", 0.10 + 0.80 * (done / max(total, 1)))

    n_done = 0
    if texts:
        results = [_maybe_pangu(t, cfg) for t in
                   translator.translate_texts(texts, glossary, tcb)]
        check_cancel()
        report("正在写回 PowerPoint…", 0.93)
        for u, tr in zip(units, results):
            if not _translated(u.text, tr, cfg):
                continue          # 模型原样退回（专名/编号等）→ 保留原文
            _set_text(u, tr)
            n_done += 1

        hits = getattr(translator, "cache_hits", 0)
        if hits:
            report(f"持久缓存命中 {hits} 段，未重复计费", 0.94)
        failed = getattr(translator, "failed_texts", 0)
        if failed:
            report(f"注意：{failed} 段因网络/服务错误未翻译，已保留原文——"
                   "重新运行即可补齐", 0.94)
    else:
        report("演示文稿中未找到可翻译的文字", 0.9)

    from .paths import atomic_output
    with atomic_output(output_path) as _out:
        prs.save(_out.tmp)        # 原子写出：中途失败不留半截损坏演示文稿
    report("完成", 1.0)
    return {"pages": len(prs.slides._sldIdLst), "blocks": n_done,
            "output": _out.path, "mode": getattr(cfg, "output_mode", "translated"),
            "backend": "pptx"}
