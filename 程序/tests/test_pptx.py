"""PowerPoint 翻译（A3）：文本框/表格/备注覆盖、样式保真、格式分派。"""
import subprocess
import sys

import pytest

from src.config import load_config
from src.pipeline import output_suffix, translate_document
from tests.conftest import ROOT

pptx = pytest.importorskip("pptx", reason="需要 python-pptx")
DECK = ROOT / "samples" / "sample_slides.pptx"


@pytest.fixture(scope="module")
def src_deck():
    if not DECK.exists():
        subprocess.run([sys.executable, "samples/make_pptx_sample.py"],
                       cwd=str(ROOT), check=True, capture_output=True)
    if not DECK.exists():
        pytest.skip("无法生成 PPTX 样张")
    return str(DECK)


@pytest.fixture(scope="module")
def translated(src_deck, tmp_path_factory):
    out = tmp_path_factory.mktemp("pptx") / "zh.pptx"
    res = translate_document(src_deck, str(out), load_config(), mock=True)
    from pptx import Presentation
    return Presentation(str(out)), res


def _cjk(s):
    return any("一" <= c <= "鿿" for c in s)


def test_returns_pptx_backend(translated):
    _, res = translated
    assert res["backend"] == "pptx" and res["blocks"] > 8 and res["pages"] == 3


def test_titles_and_body_translated(translated):
    deck, _ = translated
    texts = [sh.text_frame.text for sl in deck.slides for sh in sl.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()]
    assert texts and all(_cjk(t) for t in texts), "所有文本框应译为中文"


def test_bullet_paragraphs_preserved(translated):
    deck, _ = translated
    multi = [sh for sl in deck.slides for sh in sl.shapes
             if sh.has_text_frame and len(sh.text_frame.paragraphs) >= 4]
    assert multi, "要点页应保留 4 段结构"


def test_table_cells_translated(translated):
    deck, _ = translated
    tables = [sh.table for sl in deck.slides for sh in sl.shapes if sh.has_table]
    assert tables, "应有表格"
    cells = [c.text for t in tables for row in t.rows for c in row.cells]
    assert len(cells) == 9 and all(_cjk(c) for c in cells)


def test_notes_translated(translated):
    deck, _ = translated
    notes = [sl.notes_slide.notes_text_frame.text for sl in deck.slides
             if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip()]
    assert notes and any(_cjk(n) for n in notes), "演讲者备注应翻译"


def test_font_size_preserved(translated):
    from pptx.util import Pt
    deck, _ = translated
    tbl = [sh.table for sl in deck.slides for sh in sl.shapes if sh.has_table][0]
    assert tbl.cell(1, 1).text_frame.paragraphs[0].runs[0].font.size == Pt(14)


def test_reopenable(translated):
    """译后文件必须能被重新打开（结构未损坏）。"""
    deck, _ = translated
    assert len(deck.slides._sldIdLst) == 3


def test_pptx_suffix_naming():
    assert output_suffix("translated", ".pptx") == "_translation"


def test_legacy_ppt_rejected(tmp_path):
    from src.translator import TranslatorError
    f = tmp_path / "old.ppt"
    f.write_bytes(b"fake")
    with pytest.raises(TranslatorError, match="另存为"):
        translate_document(str(f), str(tmp_path / "o.ppt"), load_config())


def test_multilang_target_english(src_deck, tmp_path):
    """英文源、目标英语：已是英语的内容应原样保留（不无谓翻译）。"""
    out = tmp_path / "en.pptx"
    res = translate_document(src_deck, str(out),
                             load_config(source_lang="en", target_lang="en"),
                             mock=True)
    # 源本就是英文 → looks_like(en) 命中 → 0 段可译
    assert res["blocks"] == 0
